"""
Skapar Kursmaterial_Strategiportfoljen_v303.pptx
Baserad på v302 — uppdaterar versionsnummer och lägger till ny slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import shutil, copy
from lxml import etree

src = r"C:\Users\hejma\Projekt_Claude\strategiportfoljen\Kursmaterial_Strategiportfoljen_v302.pptx"
dst = r"C:\Users\hejma\Projekt_Claude\strategiportfoljen\Kursmaterial_Strategiportfoljen_v303.pptx"
shutil.copy2(src, dst)

prs = Presentation(dst)

NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xE2, 0xE8, 0xF0)
MUTED  = RGBColor(0x94, 0xA3, 0xB8)

# ── Ersätt versionssträngar i alla slides ──
replacements = [
    ('v3.02', 'v3.03'),
    ('v302',  'v303'),
    ('3.02',  '3.03'),
]

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements:
                    if old in run.text:
                        run.text = run.text.replace(old, new)

# ── Lägg till ny slide med v3.03-nyheter ──
# Kopiera sista slidens XML som bas för layout
last_slide = prs.slides[-1]
slide_layout = prs.slide_layouts[6]  # Blank layout

new_slide = prs.slides.add_slide(slide_layout)

W = prs.slide_width   # ~9144000 EMU (widescreen)
H = prs.slide_height  # ~5143500 EMU

MARGIN = Inches(0.5)

# ── Bakgrundsfärg (mörkblå, samma stil som övriga slides) ──
background = new_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = NAVY

# ── Rubrik-textruta ──
def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

# Titelrad
add_textbox(new_slide,
    left=MARGIN, top=Inches(0.3),
    width=W - 2*MARGIN, height=Inches(0.7),
    text="Nyheter i version 3.03",
    font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Underrubrik
add_textbox(new_slide,
    left=MARGIN, top=Inches(0.95),
    width=W - 2*MARGIN, height=Inches(0.4),
    text="Kontohantering, allow-list och buggfixar",
    font_size=14, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

# Punktlista
nyheter = [
    ("Allow-list-import",
     "Enbart de 6 konfigurerade kontona importeras — okänd data filtreras automatiskt."),
    ("Sparkonto alltid i Avstämning",
     "Avanza Sparande Martin visas alltid i Saldon per konto, med ledtext om saldo saknas."),
    ("Kontonummer under kontonamn",
     "Lättare att matcha appens rader mot rätt konto i Avanza Min ekonomi."),
    ("Kassa-dropdown — alltid 6 konton",
     "Alla konton visas oavsett om transaktionsfil är importerad."),
    ("Buggfix: exakt kontoidentifiering",
     "Liknamnda konton visas inte längre som duplikat — exakt matchning provas först."),
    ("iOS-kompatibla importknappar",
     "label+for-element ger pålitlig filväljare i Safari på iPad."),
]

ROW_H = Inches(0.55)
TOP_START = Inches(1.5)

for i, (rubrik, beskrivning) in enumerate(nyheter):
    top = TOP_START + i * ROW_H

    # Rubrik (fet, accent)
    txR = new_slide.shapes.add_textbox(MARGIN, top, Inches(2.4), ROW_H)
    tf = txR.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = rubrik
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    # Beskrivning (normal, ljus)
    txD = new_slide.shapes.add_textbox(Inches(3.1), top, W - Inches(3.6), ROW_H)
    tf = txD.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = beskrivning
    run.font.size = Pt(11)
    run.font.color.rgb = LIGHT

# Sidfot
add_textbox(new_slide,
    left=MARGIN, top=H - Inches(0.45),
    width=W - 2*MARGIN, height=Inches(0.35),
    text="Strategiportföljen — Kursmaterial för nybörjare · v3.03",
    font_size=9, bold=False, color=MUTED,
    align=PP_ALIGN.CENTER)

prs.save(dst)
print(f"Skapad: {dst}")
