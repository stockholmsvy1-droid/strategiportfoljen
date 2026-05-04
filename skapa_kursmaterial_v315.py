"""
Skapar Kursmaterial_Strategiportfoljen_v315.pptx
Baserad på v313, uppdaterar versionssträngar och lägger till ny slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import shutil

DIR = r"C:\Users\hejma\Projekt_Claude\strategiportfoljen"
src = fr"{DIR}\Kursmaterial_Strategiportfoljen_v313.pptx"
dst = fr"{DIR}\Kursmaterial_Strategiportfoljen_v315.pptx"
shutil.copy2(src, dst)

prs = Presentation(dst)
NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT= RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in [("v3.13","v3.15"),("v313","v315"),("3.13","3.15")]:
                    if old in run.text:
                        run.text = run.text.replace(old, new)

W, H = prs.slide_width, prs.slide_height
slide_layout = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(slide_layout)

bg = new_slide.background; fill = bg.fill
fill.solid(); fill.fore_color.rgb = NAVY

MARGIN = Inches(0.5)

def add_tb(slide, left, top, width, height, text, size, bold=False,
           color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

add_tb(new_slide, MARGIN, Inches(0.25), W-2*MARGIN, Inches(0.7),
       "Nyheter i version 3.15", 26, bold=True, color=WHITE)
add_tb(new_slide, MARGIN, Inches(0.9), W-2*MARGIN, Inches(0.4),
       "Ny Översikt-flik — cockpit med signaler, nyckeltal och portfölj-accordion",
       13, bold=False, color=ACCENT)

nyheter = [
    ("Ny Översikt-flik",        "Startsida ombyggd till cockpit — period-väljare (1D–Allt) styr hela appen"),
    ("Cockpit: 4 nyckeltal",    "Portföljvärde, periodsavkastning, nettoinsatt kapital och tillgänglig likviditet"),
    ("Portfölj-accordion",      "Alla innehav per kategori, fällbara sektioner, inline MA200-redigering direkt i listan"),
    ("Signal-pills",            "Grön/gul/röd plupp med tooltip per innehav — följer tvådagarsregeln för kat. 3–6"),
    ("Kategori-donut",          "Fördelningsdiagram med signalcirkel, pil och viktbalans-signal i legenden"),
    ("Navigation 5 flikar",     "Översikt · Portfölj · Signaler · Kategorier · Mer — tydligare struktur"),
    ("Fix: interna överföringar","Kontoöverföringar mellan egna Avanza-konton räknas inte som insatt kapital"),
]

ROW_H = Inches(0.48)
TOP_S = Inches(1.42)
for i, (rubrik, beskr) in enumerate(nyheter):
    top = TOP_S + i * ROW_H
    add_tb(new_slide, MARGIN, top, Inches(2.5), ROW_H,
           rubrik, 10, bold=True, color=ACCENT)
    add_tb(new_slide, Inches(3.1), top, W-Inches(3.6), ROW_H,
           beskr, 10, bold=False, color=LIGHT)

add_tb(new_slide, MARGIN, H-Inches(0.45), W-2*MARGIN, Inches(0.35),
       "Strategiportföljen — Kursmaterial · v3.15", 9, color=MUTED,
       align=PP_ALIGN.CENTER)

prs.save(dst)
print(f"Skapad: {dst}")
