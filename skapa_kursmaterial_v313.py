"""
Skapar Kursmaterial_Strategiportfoljen_v313.pptx
Baserad på v310, uppdaterar versionssträngar och lägger till ny slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import shutil

DIR = r"C:\Users\hejma\Projekt_Claude\strategiportfoljen"
src = fr"{DIR}\Kursmaterial_Strategiportfoljen_v310.pptx"
dst = fr"{DIR}\Kursmaterial_Strategiportfoljen_v313.pptx"
shutil.copy2(src, dst)

prs = Presentation(dst)
NAVY  = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT= RGBColor(0x0E, 0xA5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in [("v3.10","v3.13"),("v310","v313"),("3.10","3.13")]:
                    if old in run.text:
                        run.text = run.text.replace(old, new)

W, H = prs.slide_width, prs.slide_height
slide_layout = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(slide_layout)

bg = new_slide.background; fill = bg.fill
fill.solid(); fill.fore_color.rgb = NAVY

MARGIN = Inches(0.5)

def add_tb(slide, left, top, width, height, text, size, bold=False,
           color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

add_tb(new_slide, MARGIN, Inches(0.25), W-2*MARGIN, Inches(0.7),
       "Nyheter i version 3.13", 26, bold=True, color=WHITE)
add_tb(new_slide, MARGIN, Inches(0.9), W-2*MARGIN, Inches(0.4),
       "Inställningar-sektion och förbättrad Kassa",
       13, bold=False, color=ACCENT)

nyheter = [
    ("Kontokonfiguration",    "Lägg till/redigera/ta bort Avanza-konton från UI — ingen kodredigering"),
    ("Kategori-editor",       "Ersätter prompt()-dialoger med visuellt inline-formulär (emoji, färg, vikter, signal)"),
    ("Strategiparametrar",    "MA200-gränser, nödutgång, ombalansering och koncentrationsrisk — konfigurerbara"),
    ("Profil & information",  "Redigerbara fält för namn och strategi-titel som används i exporter"),
    ("Värdepappersfilter",    "Exkludera VP/konton vid import — redigerbara listor i UI"),
    ("Export/import JSON",    "Flytta hela strategikonfigurationen mellan enheter som JSON-fil"),
    ("Kassa inline-inmatning","Alla konton alltid synliga, inmatning direkt per rad — inget separat formulär"),
]

ROW_H = Inches(0.48)
TOP_S = Inches(1.42)
for i, (rubrik, beskr) in enumerate(nyheter):
    top = TOP_S + i * ROW_H
    add_tb(new_slide, MARGIN, top, Inches(2.5), ROW_H,
           rubrik, 10, bold=True, color=ACCENT)
    add_tb(new_slide, Inches(3.1), top, W-Inches(3.6), ROW_H,
           beskr, 10, bold=False, color=LIGHT)

add_tb(new_slide, MARGIN, H-Inches(0.45), W-2*MARGIN, Inches(0.35),
       "Strategiportföljen — Kursmaterial · v3.13", 9, color=MUTED,
       align=PP_ALIGN.CENTER)

prs.save(dst)
print(f"Skapad: {dst}")
